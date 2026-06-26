"""PPTX pipeline (.pptx via python-pptx).

Extracts text from slides in slide order, including text boxes, tables,
grouped shapes, and speaker notes when present. The rendered text is handed
to the shared text indexer with slide headings.

read_segment supports slide_start / slide_end ranges, page_start / page_end
aliases for slide ranges, regex pattern search, and generic offset/max_chars
chunking over the rendered deck text.

Legacy binary .ppt is not supported; users should resave to .pptx or PDF.
"""
from __future__ import annotations

import io
import logging
import re
from typing import Any

from marginalia.pipelines._text_indexer import index_extracted_text
from marginalia.pipelines.base import (
    Pipeline,
    PipelineContext,
    PipelineResult,
    SegmentResult,
)
from marginalia.pipelines.registry import register_pipeline
from marginalia.storage.base import StorageBackend

log = logging.getLogger(__name__)

# Do not reject PPTX by compressed package size. Slide/text extraction and the
# prompt text cap below control indexing cost more directly than package bytes.
MAX_PPTX_SLIDES = 400
MAX_OUTPUT_CHARS = 100_000
DEFAULT_MAX_CHARS = 8000


@register_pipeline(
    mimes=(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint.presentation.macroenabled.12",
    ),
    exts=(".pptx", ".pptm"),
)
class PptxPipeline(Pipeline):
    name = "pptx"

    async def run(
        self,
        *,
        ctx: PipelineContext,
        storage: StorageBackend,
    ) -> PipelineResult:
        slides, coverage = await self._extract_slides_with_coverage(
            storage, ctx.storage_key, max_slides=MAX_PPTX_SLIDES,
        )
        full_body = "\n\n".join(slides)
        indexed_chars = min(len(full_body), MAX_OUTPUT_CHARS)
        body = full_body
        partial_reasons = list(coverage.get("partial_reasons") or [])
        if len(full_body) > MAX_OUTPUT_CHARS:
            body = full_body[:MAX_OUTPUT_CHARS] + "\n[...deck truncated for indexing...]"
            partial_reasons.append("prompt_text_cap")
        coverage.update({
            "total_chars": len(full_body),
            "indexed_chars": indexed_chars,
            "indexed_partial": bool(partial_reasons),
            "partial_reasons": list(dict.fromkeys(partial_reasons)),
            "max_index_chars": MAX_OUTPUT_CHARS,
            "max_index_slides": MAX_PPTX_SLIDES,
            "text_truncated": bool(partial_reasons),
        })
        return await index_extracted_text(
            body, ctx, kind="text", coverage=coverage,
        )

    async def read_segment(
        self,
        *,
        file_row: Any,
        args: dict[str, Any],
        storage: StorageBackend,
    ) -> SegmentResult:
        slides, _coverage = await self._extract_slides_with_coverage(
            storage, file_row.storage_key,
        )
        return self._slice(slides, args, file_row=file_row)

    async def read_segment_from_bytes(
        self,
        body: bytes,
        args: dict[str, Any],
        *,
        filename: str | None = None,
    ) -> SegmentResult:
        """Bytes-first variant used by ArchivePipeline for member peeks."""
        try:
            slides, _coverage = self._render_from_bytes_with_coverage(body)
        except Exception as exc:  # noqa: BLE001
            return SegmentResult(error=f"pptx parse failed: {exc}")
        return self._slice(slides, args, file_row=None)

    def _slice(
        self,
        slides: list[str],
        args: dict[str, Any],
        *,
        file_row: Any | None,
    ) -> SegmentResult:
        body = "\n\n".join(slides)
        total_slides = len(slides)

        offset = max(0, int(args.get("offset") or 0))
        max_chars = int(args.get("max_chars") or DEFAULT_MAX_CHARS)
        if max_chars <= 0:
            max_chars = DEFAULT_MAX_CHARS

        pattern = (args.get("pattern") or "").strip()
        scoped = _resolve_slide_window(args, total_slides=total_slides)
        if isinstance(scoped, SegmentResult):
            return scoped

        if pattern:
            scope_slides = slides
            slide_offset = 0
            if scoped is not None:
                start, end = scoped
                scope_slides = slides[start - 1:end]
                slide_offset = start - 1
            return _pptx_pattern_search(
                slides=scope_slides,
                pattern=pattern,
                context_lines=int(args.get("context_lines") or 2),
                max_matches=int(args.get("max_matches") or 20),
                match_offset=max(0, int(args.get("match_offset") or 0)),
                slide_offset=slide_offset,
                total_slides_full=total_slides,
            )

        if scoped is not None:
            start, end = scoped
            slab = "\n\n".join(slides[start - 1:end])
            return _clamp_pptx(
                slab,
                offset,
                max_chars,
                extras={
                    "slide_start": start,
                    "slide_end": end,
                    "total_slides": total_slides,
                },
            )

        if any(args.get(key) for key in ("section_id", "heading", "line_start", "line_end")):
            from marginalia.pipelines.text import TextPipeline

            return TextPipeline()._slice(body=body, args=args, file_row=file_row)

        slide_start = body[:offset].count("\n# Slide ") + 1
        chunk = body[offset: offset + max_chars]
        slide_end = min(total_slides, slide_start + chunk.count("\n# Slide "))
        return _clamp_pptx(
            body,
            offset,
            max_chars,
            extras={
                "slide_start": slide_start,
                "slide_end": slide_end,
                "total_slides": total_slides,
            },
        )

    @classmethod
    async def _extract_slides_with_coverage(
        cls,
        storage: StorageBackend,
        key: str,
        *,
        max_slides: int | None = None,
    ) -> tuple[list[str], dict[str, Any]]:
        try:
            from pptx import Presentation  # type: ignore  # noqa: F401
        except ImportError as exc:
            raise RuntimeError(
                "pptx pipeline needs python-pptx; `pip install python-pptx`"
            ) from exc

        buf = bytearray()
        async for chunk in storage.get(key):
            buf.extend(chunk)
        slides, coverage = cls._render_from_bytes_with_coverage(
            bytes(buf), max_slides=max_slides,
        )
        coverage["total_bytes"] = len(buf)
        coverage["indexed_bytes"] = len(buf)
        return slides, coverage

    @staticmethod
    def _render_from_bytes_with_coverage(
        body: bytes,
        *,
        max_slides: int | None = None,
    ) -> tuple[list[str], dict[str, Any]]:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError as exc:
            raise RuntimeError(
                "pptx pipeline needs python-pptx; `pip install python-pptx`"
            ) from exc

        prs = Presentation(io.BytesIO(body))
        total_slides = len(prs.slides)
        slide_limit = (
            max(0, int(max_slides))
            if max_slides is not None and int(max_slides) > 0
            else None
        )
        slides: list[str] = []
        shape_count = 0
        table_count = 0
        notes_count = 0
        for idx, slide in enumerate(prs.slides, start=1):
            if slide_limit is not None and len(slides) >= slide_limit:
                break
            rendered, stats = _render_slide(slide, idx)
            slides.append(rendered)
            shape_count += stats["shape_count"]
            table_count += stats["table_count"]
            notes_count += stats["notes_count"]
        indexed_partial = len(slides) < total_slides
        coverage = {
            "unit": "slides",
            "source_mode": "pptx_extracted_text",
            "total_units": total_slides,
            "indexed_units": len(slides),
            "total_slides": total_slides,
            "indexed_slides": len(slides),
            "shape_count": shape_count,
            "table_count": table_count,
            "notes_count": notes_count,
            "indexed_partial": indexed_partial,
            "partial_reasons": ["slide_cap"] if indexed_partial else [],
            "chunked": False,
            "chunk_count": 1,
            "text_truncated": indexed_partial,
        }
        if slide_limit is not None:
            coverage["max_index_slides"] = slide_limit
        return slides, coverage


def _render_slide(slide: Any, slide_no: int) -> tuple[str, dict[str, int]]:
    title = _shape_plain_text(getattr(slide.shapes, "title", None))
    title_line = f"# Slide {slide_no}" + (f": {title}" if title else "")
    lines = [title_line]
    shape_count = 0
    table_count = 0

    for line, stats in _iter_shape_lines(slide.shapes, skip_shape=slide.shapes.title):
        if line:
            lines.append(line)
        shape_count += stats["shape_count"]
        table_count += stats["table_count"]

    notes = _notes_text(slide)
    notes_count = 1 if notes else 0
    if notes:
        lines.append("Notes:")
        lines.extend(notes)

    if len(lines) == 1:
        lines.append("(no extractable text)")
    return "\n".join(lines), {
        "shape_count": shape_count,
        "table_count": table_count,
        "notes_count": notes_count,
    }


def _iter_shape_lines(shapes: Any, *, skip_shape: Any = None):
    for shape in shapes:
        if skip_shape is not None and shape == skip_shape:
            continue
        shape_count = 1
        table_count = 0
        if hasattr(shape, "shapes"):
            for line, stats in _iter_shape_lines(shape.shapes):
                yield line, stats
            continue
        if getattr(shape, "has_table", False):
            for row_idx, row in enumerate(shape.table.rows):
                cells = [
                    _clean_cell_text(cell.text)
                    for cell in row.cells
                ]
                yield " | ".join(cells), {
                    "shape_count": shape_count if row_idx == 0 else 0,
                    "table_count": 1 if row_idx == 0 else 0,
                }
            continue
        text = _shape_plain_text(shape)
        if text:
            for line_idx, line in enumerate(text.splitlines()):
                cleaned = line.strip()
                if cleaned:
                    yield cleaned, {
                        "shape_count": shape_count if line_idx == 0 else 0,
                        "table_count": table_count,
                    }


def _shape_plain_text(shape: Any) -> str:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return ""
    lines: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if not text:
            text = (paragraph.text or "").strip()
        if text:
            lines.append(text)
    return "\n".join(lines).strip()


def _clean_cell_text(text: str) -> str:
    return (text or "").strip().replace("\n", " ").replace("|", r"\|")


def _notes_text(slide: Any) -> list[str]:
    try:
        if not slide.has_notes_slide:
            return []
        frame = slide.notes_slide.notes_text_frame
    except Exception:  # noqa: BLE001
        return []
    lines: list[str] = []
    for paragraph in frame.paragraphs:
        text = (paragraph.text or "").strip()
        if text and text.lower() != "click to add notes":
            lines.append(text)
    return lines


def _resolve_slide_window(
    args: dict[str, Any],
    *,
    total_slides: int,
) -> tuple[int, int] | SegmentResult | None:
    start_raw = args.get("slide_start", args.get("page_start"))
    end_raw = args.get("slide_end", args.get("page_end"))
    if start_raw in (None, "") and end_raw in (None, ""):
        return None
    if total_slides <= 0:
        return SegmentResult(error="pptx has no slides")
    try:
        start = max(1, int(start_raw)) if start_raw not in (None, "") else 1
        end = int(end_raw) if end_raw not in (None, "") else start
    except (TypeError, ValueError):
        return SegmentResult(error="slide_start/slide_end must be integers")
    if end < start:
        return SegmentResult(error="slide_end must be >= slide_start")
    start = max(1, min(start, total_slides))
    end = max(start, min(end, total_slides))
    return start, end


def _clamp_pptx(
    text: str,
    offset: int,
    max_chars: int,
    *,
    extras: dict[str, Any] | None = None,
) -> SegmentResult:
    extras = dict(extras or {})
    total = len(text)
    chunk = text[offset: offset + max_chars]
    truncated = (offset + len(chunk)) < total
    extras.update({
        "offset": offset,
        "char_count": len(chunk),
        "total_chars": total,
        "truncated": truncated,
    })
    if truncated:
        extras["next_offset"] = offset + len(chunk)
    if not chunk:
        return SegmentResult(text="", error="empty result", extras=extras)
    return SegmentResult(text=chunk, extras=extras)


def _pptx_pattern_search(
    *,
    slides: list[str],
    pattern: str,
    context_lines: int,
    max_matches: int,
    match_offset: int = 0,
    slide_offset: int = 0,
    total_slides_full: int | None = None,
) -> SegmentResult:
    try:
        rx = re.compile(pattern, re.IGNORECASE | re.MULTILINE)
    except re.error as exc:
        return SegmentResult(error=f"invalid regex: {exc}")

    full_total = total_slides_full if total_slides_full is not None else len(slides)
    all_hits: list[dict[str, Any]] = []
    for slide_idx, slide_text in enumerate(slides, start=1):
        if not slide_text:
            continue
        slide_no = slide_idx + slide_offset
        lines = slide_text.splitlines()
        for line_idx, line in enumerate(lines, start=1):
            for match in rx.finditer(line):
                start = max(0, line_idx - 1 - context_lines)
                end = min(len(lines), line_idx + context_lines)
                all_hits.append({
                    "slide": slide_no,
                    "line": line_idx,
                    "match": match.group(0)[:200],
                    "context": "\n".join(lines[start:end]),
                })

    total = len(all_hits)
    hits = all_hits[match_offset: match_offset + max_matches]
    has_more = (match_offset + len(hits)) < total
    extras: dict[str, Any] = {
        "pattern": pattern,
        "match_count": len(hits),
        "total_matches": total,
        "match_offset": match_offset,
        "has_more": has_more,
        "hits": hits,
        "total_slides": full_total,
    }
    if slide_offset:
        extras["scope_slide_start"] = slide_offset + 1
        extras["scope_slide_end"] = slide_offset + len(slides)
    if has_more:
        extras["next_match_offset"] = match_offset + len(hits)

    if not hits:
        if match_offset and total:
            err = f"match_offset {match_offset} exceeds total_matches {total}"
        else:
            err = "no matches"
        return SegmentResult(text="", error=err, extras=extras)

    rendered = "\n\n".join(
        f"[Slide {h['slide']} L{h['line']}] {h['match']}\n  > {h['context']}"
        for h in hits
    )
    return SegmentResult(text=rendered, extras=extras)
