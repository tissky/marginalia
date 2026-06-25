from __future__ import annotations

import pytest

from marginalia.pipelines.base import PipelineContext
from marginalia.pipelines.docx import DocxPipeline
from marginalia.pipelines import pptx as pptx_module
from marginalia.pipelines.pptx import PptxPipeline
from marginalia.pipelines.spreadsheet import SpreadsheetPipeline


OLD_DOCX_XLSX_CAP = 30 * 1024 * 1024
OLD_PPTX_CAP = 50 * 1024 * 1024


class _ChunkStorage:
    def __init__(self, size: int):
        self.size = size

    async def get(self, key: str):
        del key
        yield b"x" * self.size


def _ctx(name: str, size: int) -> PipelineContext:
    return PipelineContext(
        file_id=f"file-{name}",
        storage_key=f"{name}.bin",
        sha256="0" * 64,
        size_bytes=size,
        mime_type=None,
        original_ext=f".{name}",
        folder_path="/tests",
        sibling_names=[],
        display_name=f"{name}.bin",
    )


@pytest.mark.asyncio
async def test_docx_ingest_no_longer_rejects_by_package_size(monkeypatch):
    seen: dict[str, int] = {}
    size = OLD_DOCX_XLSX_CAP + 1

    def fake_parse(body: bytes) -> list[str]:
        seen["bytes"] = len(body)
        return ["Large docx parsed"]

    async def fake_index(body, ctx, *, kind, coverage):
        return {
            "body": body,
            "ctx": ctx,
            "kind": kind,
            "coverage": coverage,
        }

    monkeypatch.setattr(DocxPipeline, "_parse_paragraphs_from_bytes", staticmethod(fake_parse))
    monkeypatch.setattr("marginalia.pipelines.docx.index_extracted_text", fake_index)

    result = await DocxPipeline().run(
        ctx=_ctx("docx", size),
        storage=_ChunkStorage(size),
    )

    assert seen["bytes"] == size
    assert result["coverage"]["total_paragraphs"] == 1


@pytest.mark.asyncio
async def test_spreadsheet_ingest_no_longer_rejects_by_package_size(monkeypatch):
    seen: dict[str, int] = {}
    size = OLD_DOCX_XLSX_CAP + 1

    def fake_render(body: bytes):
        seen["bytes"] = len(body)
        return "# Sheet: large\nvalue", {
            "unit": "rows",
            "source_mode": "spreadsheet_row_sample",
            "total_rows": 1,
            "indexed_rows": 1,
        }

    async def fake_index(body, ctx, *, kind, coverage):
        return {
            "body": body,
            "ctx": ctx,
            "kind": kind,
            "coverage": coverage,
        }

    monkeypatch.setattr(
        SpreadsheetPipeline,
        "_render_from_bytes_with_coverage",
        staticmethod(fake_render),
    )
    monkeypatch.setattr("marginalia.pipelines.spreadsheet.index_extracted_text", fake_index)

    result = await SpreadsheetPipeline().run(
        ctx=_ctx("xlsx", size),
        storage=_ChunkStorage(size),
    )

    assert seen["bytes"] == size
    assert result["coverage"]["total_bytes"] == size


@pytest.mark.asyncio
async def test_pptx_ingest_no_longer_rejects_by_package_size(monkeypatch):
    seen: dict[str, int] = {}
    size = OLD_PPTX_CAP + 1

    def fake_render(body: bytes, *, max_slides=None):
        assert max_slides == pptx_module.MAX_PPTX_SLIDES
        seen["bytes"] = len(body)
        return ["# Slide 1\nLarge pptx parsed"], {
            "unit": "slides",
            "source_mode": "pptx_extracted_text",
            "total_slides": 1,
            "indexed_slides": 1,
        }

    async def fake_index(body, ctx, *, kind, coverage):
        return {
            "body": body,
            "ctx": ctx,
            "kind": kind,
            "coverage": coverage,
        }

    monkeypatch.setattr(PptxPipeline, "_render_from_bytes_with_coverage", staticmethod(fake_render))
    monkeypatch.setattr("marginalia.pipelines.pptx.index_extracted_text", fake_index)

    result = await PptxPipeline().run(
        ctx=_ctx("pptx", size),
        storage=_ChunkStorage(size),
    )

    assert seen["bytes"] == size
    assert result["coverage"]["total_slides"] == 1


def _build_pptx(slide_count: int) -> bytes:
    from io import BytesIO

    from pptx import Presentation  # type: ignore

    prs = Presentation()
    for idx in range(1, slide_count + 1):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Slide {idx}"
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_ingest_is_limited_by_slide_budget(monkeypatch):
    monkeypatch.setattr(pptx_module, "MAX_PPTX_SLIDES", 2)

    slides, coverage = PptxPipeline._render_from_bytes_with_coverage(
        _build_pptx(3),
        max_slides=pptx_module.MAX_PPTX_SLIDES,
    )

    assert len(slides) == 2
    assert coverage["total_slides"] == 3
    assert coverage["indexed_slides"] == 2
    assert coverage["indexed_partial"] is True
    assert coverage["partial_reasons"] == ["slide_cap"]
