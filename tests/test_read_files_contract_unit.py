from __future__ import annotations

import io
import zipfile
from types import SimpleNamespace

import pytest

from marginalia.pipelines.archive import ArchivePipeline
from marginalia.pipelines.docx import DocxPipeline
from marginalia.pipelines.markitdown import MarkItDownPipeline
from marginalia.pipelines.pptx import PptxPipeline
from marginalia.pipelines.spreadsheet import SpreadsheetPipeline


class _MemoryStorage:
    def __init__(self, body: bytes) -> None:
        self.body = body

    async def get(self, key: str):
        assert key == "file-key"
        yield self.body


def _xlsx_with_late_needle() -> bytes:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "data"
    ws.append(["row", "value"])
    for idx in range(1, 260):
        ws.append([idx, f"value-{idx}"])
    ws.append([260, "needle-after-row-cap"])
    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


def _docx_with_heading() -> bytes:
    from docx import Document

    doc = Document()
    doc.add_heading("Contract Heading", level=1)
    doc.add_paragraph("Heading body line")
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def _pptx_with_heading() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Contract Slide"
    box = slide.shapes.add_textbox(914400, 1371600, 5486400, 914400)
    box.text = "Slide body line"
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


@pytest.mark.asyncio
async def test_markitdown_heading_read_uses_full_extracted_text(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import marginalia.pipelines.markitdown as mod

    monkeypatch.setattr(
        mod,
        "_convert_bytes_with_markitdown",
        lambda body, suffix: "A" * 125_000 + "\n# Late Chapter\nneedle-after-index-cap",
    )

    result = await MarkItDownPipeline().read_segment_from_bytes(
        b"source bytes",
        {"heading": "Late Chapter", "max_chars": 200},
        filename="book.epub",
    )

    assert result.error is None
    assert "needle-after-index-cap" in result.text
    assert result.extras["located_via"] == "body-heading-scan"


@pytest.mark.asyncio
async def test_spreadsheet_read_uses_full_rows_not_ingest_sample() -> None:
    result = await SpreadsheetPipeline().read_segment(
        file_row=SimpleNamespace(storage_key="file-key"),
        args={"pattern": "needle-after-row-cap", "context_lines": 0},
        storage=_MemoryStorage(_xlsx_with_late_needle()),
    )

    assert result.error is None
    assert "needle-after-row-cap" in result.text
    assert result.extras["total_matches"] == 1


@pytest.mark.asyncio
async def test_docx_heading_and_line_reads_use_complete_extracted_text() -> None:
    pipeline = DocxPipeline()
    row = SimpleNamespace(storage_key="file-key", description=None)

    heading = await pipeline.read_segment(
        file_row=row,
        args={"heading": "Contract Heading", "max_chars": 200},
        storage=_MemoryStorage(_docx_with_heading()),
    )
    assert heading.error is None
    assert "Heading body line" in heading.text

    line = await pipeline.read_segment(
        file_row=row,
        args={"line_start": 1, "line_end": 2},
        storage=_MemoryStorage(_docx_with_heading()),
    )
    assert line.error is None
    assert "Contract Heading" in line.text
    assert line.extras["line_start"] == 1


@pytest.mark.asyncio
async def test_pptx_heading_and_line_reads_use_complete_extracted_text() -> None:
    pipeline = PptxPipeline()
    row = SimpleNamespace(storage_key="file-key", description=None)

    heading = await pipeline.read_segment(
        file_row=row,
        args={"heading": "Slide 1: Contract Slide", "max_chars": 300},
        storage=_MemoryStorage(_pptx_with_heading()),
    )
    assert heading.error is None
    assert "Slide body line" in heading.text

    line = await pipeline.read_segment(
        file_row=row,
        args={"line_start": 1, "line_end": 2},
        storage=_MemoryStorage(_pptx_with_heading()),
    )
    assert line.error is None
    assert "Contract Slide" in line.text
    assert line.extras["line_start"] == 1


@pytest.mark.asyncio
async def test_archive_member_read_dispatches_to_full_inner_reader() -> None:
    archive_body = io.BytesIO()
    with zipfile.ZipFile(archive_body, "w") as zf:
        zf.writestr("reports/data.xlsx", _xlsx_with_late_needle())

    result = await ArchivePipeline().read_segment(
        file_row=SimpleNamespace(storage_key="file-key", original_ext=".zip", id="file-id"),
        args={
            "member_path": "reports/data.xlsx",
            "pattern": "needle-after-row-cap",
            "context_lines": 0,
        },
        storage=_MemoryStorage(archive_body.getvalue()),
    )

    assert result.error is None
    assert "needle-after-row-cap" in result.text
    assert result.extras["total_matches"] == 1
