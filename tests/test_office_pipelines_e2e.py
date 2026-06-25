"""End-to-end tests for docx + pptx + spreadsheet pipelines.

Run:
    .venv/Scripts/python tests/test_office_pipelines_e2e.py

Verifies for both pipelines:
  1. Registry routes the right mime/ext to the right pipeline.
  2. Round-trip: real .docx / .pptx / .xlsx file is built in-memory, ingested,
     fake LLM returns a valid index, file row reaches ingest_status='done'.
  3. read_segment supports its specialty fields:
       docx        → paragraph_start / paragraph_end, pattern, chunk
       pptx        → slide_start / slide_end, pattern, chunk
       spreadsheet → heading="Sheet: <name>", pattern, chunk
  4. Invalid args (out-of-range paragraph, unknown sheet) return ok=false
     with a clear error message.

No real LLM is called — `get_chat_client("ingest")` is replaced with a
fake that returns a valid PipelineResult-shaped JSON.
"""
from __future__ import annotations

import os
from uuid import uuid4
import asyncio
import io
import sys
from pathlib import Path

_TEST_PARENT = Path(os.environ.get("MARGINALIA_TEST_TMP", Path(__file__).resolve().parent))
_TEST_ROOT = _TEST_PARENT / f"_office_pipelines_e2e_data_{os.getpid()}_{uuid4().hex[:8]}"
_TEST_ROOT.mkdir(parents=True)
os.environ["MARGINALIA_HOME"] = str(_TEST_ROOT)
os.environ["STORAGE_BACKEND"] = "local"
os.environ["WORKER_ENABLED"] = "false"
os.environ["LLM_DEFAULT_API_KEY"] = "sk-fake"
os.environ["LLM_DEFAULT_MODEL"] = "fake-model"

from sqlalchemy import select

from marginalia.config import get_settings
get_settings.cache_clear()  # type: ignore[attr-defined]

from marginalia import llm
from marginalia.db.engine import get_engine, get_session_factory
from marginalia.db.models import Base, File, FileEntry
from marginalia.llm.types import ChatRequest, ChatResponse, TokenUsage
from marginalia.pipelines.docx import DocxPipeline
from marginalia.pipelines.pptx import PptxPipeline
from marginalia.pipelines.registry import resolve_pipeline
from marginalia.pipelines.spreadsheet import SpreadsheetPipeline
from marginalia.storage import get_storage


# ---- fake LLM ---------------------------------------------------------------

CALL_LOG: list[tuple[str, ChatRequest]] = []


def _make_fake(kind: str):
    tagged = f"""<summary>
A test {kind} used to exercise the {kind} pipeline path.
</summary>
<description>
Synthetic {kind} fixture for the office pipeline.
</description>
<sections>
s1 | lines 1-3 | Overview | Top of the document. | test, fixture, {kind}
</sections>
<extra>
notable: Synthetic {kind} fixture with no real semantic content.
</extra>
<entry_extra>
Test fixture stored in /tests/office.
</entry_extra>
<catalog_path>Tests / Office</catalog_path>
<tags>
source: test-fixture
form: {kind}
</tags>"""

    class _Fake:
        profile_name = "ingest"
        model = "fake-ingest"

        async def complete(self, request: ChatRequest) -> ChatResponse:
            CALL_LOG.append((kind, request))
            return ChatResponse(
                text=tagged,
                tool_calls=[],
                stop_reason="end_turn",
                usage=TokenUsage(input_tokens=500, output_tokens=200),
                parsed_json=None,
            )

    return _Fake()


def _install_fakes() -> None:
    llm.reset_clients_cache()
    docx_fake = _make_fake("docx")
    pptx_fake = _make_fake("pptx")
    spreadsheet_fake = _make_fake("spreadsheet")
    import marginalia.pipelines._text_indexer as imod

    # Both pipelines route through index_extracted_text → get_chat_client.
    # We pick which fake to use by walking the call stack — index_extracted_text
    # is invoked from DocxPipeline.run or SpreadsheetPipeline.run, so the
    # nearest pipeline module name decides.
    def _factory(profile: str = "ingest"):
        import inspect
        for frame in inspect.stack():
            mod = frame.frame.f_globals.get("__name__", "")
            if mod.endswith(".docx"):
                return docx_fake
            if mod.endswith(".pptx"):
                return pptx_fake
            if mod.endswith(".spreadsheet"):
                return spreadsheet_fake
        return docx_fake  # safe default for tests (any will satisfy schema)
    imod.get_chat_client = _factory  # type: ignore[assignment]


# ---- fixtures ---------------------------------------------------------------

def _build_docx() -> bytes:
    """Build a real .docx with a heading, body paragraphs, and a table."""
    from docx import Document  # type: ignore

    doc = Document()
    doc.add_heading("Marginalia Test Document", level=1)
    doc.add_paragraph(
        "This is the introduction paragraph. It mentions raft consensus."
    )
    doc.add_paragraph("A second paragraph with no special keywords.")
    doc.add_heading("Algorithm Details", level=2)
    doc.add_paragraph("Leader election proceeds in three steps.")
    doc.add_paragraph("Each follower votes for at most one candidate per term.")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Term"
    table.rows[0].cells[1].text = "Leader"
    table.rows[1].cells[0].text = "1"
    table.rows[1].cells[1].text = "alice"

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _build_xlsx() -> bytes:
    """Build a real .xlsx with two sheets."""
    import openpyxl  # type: ignore

    wb = openpyxl.Workbook()
    sheet1 = wb.active
    sheet1.title = "consensus"
    sheet1.append(["term", "leader", "votes"])
    sheet1.append([1, "alice", 3])
    sheet1.append([2, "bob", 5])

    sheet2 = wb.create_sheet(title="latency")
    sheet2.append(["round", "ms", "notes"])
    sheet2.append([1, 12, "warm-up"])
    sheet2.append([2, 8, "steady"])

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _build_pptx() -> bytes:
    """Build a real .pptx with text, a table, and speaker notes."""
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches  # type: ignore

    prs = Presentation()

    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "Consensus Overview"
    slide1.placeholders[1].text = "Raft leader election\nPaxos acceptors"
    slide1.notes_slide.notes_text_frame.text = "Presenter note mentions quorum."

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Latency Table"
    table = slide2.shapes.add_table(
        2, 2, Inches(1), Inches(1.4), Inches(5), Inches(1),
    ).table
    table.cell(0, 0).text = "Round"
    table.cell(0, 1).text = "Latency"
    table.cell(1, 0).text = "steady"
    table.cell(1, 1).text = "8 ms"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---- helpers ---------------------------------------------------------------

async def _create_schema():
    engine = get_engine()
    async with engine.begin() as conn:
        await conn.run_sync(Base.metadata.create_all)


async def _seed_file(
    *, body: bytes, mime: str, name: str,
) -> tuple[str, str]:
    """Use the real upload service so Folder + FileEntry get wired correctly."""
    from marginalia.services.upload import upload

    storage = get_storage()

    async def _stream():
        yield body

    factory = get_session_factory()
    async with factory() as db:
        result = await upload(
            db, storage,
            stream=_stream(), fallback_name=name,
            remote_path=f"/tests/{name}",
            content_type=mime,
        )
        await db.commit()
        return result.file_id, result.entry_id


async def _ingest(file_id: str) -> None:
    """Run the file through ingest_file handler."""
    from marginalia.tasks.handlers.ingest_file import handle_ingest_file
    await handle_ingest_file({"file_id": file_id})


# ---- main ------------------------------------------------------------------

async def go() -> None:
    await _create_schema()
    _install_fakes()

    # 1. registry routing
    docx_pl = resolve_pipeline(
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".docx",
    )
    xlsx_pl = resolve_pipeline(
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xlsx",
    )
    pptx_pl = resolve_pipeline(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".pptx",
    )
    assert isinstance(docx_pl, DocxPipeline), docx_pl
    assert isinstance(pptx_pl, PptxPipeline), pptx_pl
    assert isinstance(xlsx_pl, SpreadsheetPipeline), xlsx_pl
    print("[1] registry routes docx→DocxPipeline, pptx→PptxPipeline, xlsx→SpreadsheetPipeline")

    # 2. round-trip ingest for all office formats
    docx_bytes = _build_docx()
    pptx_bytes = _build_pptx()
    xlsx_bytes = _build_xlsx()

    docx_file_id, _ = await _seed_file(
        body=docx_bytes,
        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        name="test.docx",
    )
    xlsx_file_id, _ = await _seed_file(
        body=xlsx_bytes,
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        name="test.xlsx",
    )
    pptx_file_id, _ = await _seed_file(
        body=pptx_bytes,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        name="test.pptx",
    )
    await _ingest(docx_file_id)
    await _ingest(pptx_file_id)
    await _ingest(xlsx_file_id)

    factory = get_session_factory()
    async with factory() as db:
        docx_row = (await db.execute(
            select(File).where(File.id == docx_file_id)
        )).scalar_one()
        xlsx_row = (await db.execute(
            select(File).where(File.id == xlsx_file_id)
        )).scalar_one()
        pptx_row = (await db.execute(
            select(File).where(File.id == pptx_file_id)
        )).scalar_one()

    assert docx_row.ingest_status == "done", docx_row.ingest_status
    assert pptx_row.ingest_status == "done", pptx_row.ingest_status
    assert xlsx_row.ingest_status == "done", xlsx_row.ingest_status
    assert docx_row.kind == "text"
    assert pptx_row.kind == "text"
    assert xlsx_row.kind == "table"
    docx_coverage = (docx_row.description or {}).get("coverage") or {}
    pptx_coverage = (pptx_row.description or {}).get("coverage") or {}
    xlsx_coverage = (xlsx_row.description or {}).get("coverage") or {}
    assert docx_coverage.get("source_mode") == "docx_extracted_text", docx_coverage
    assert docx_coverage.get("indexed_partial") is False, docx_coverage
    assert pptx_coverage.get("source_mode") == "pptx_extracted_text", pptx_coverage
    assert pptx_coverage.get("total_slides") == 2, pptx_coverage
    assert xlsx_coverage.get("source_mode") == "spreadsheet_row_sample", xlsx_coverage
    assert xlsx_coverage.get("indexed_partial") is False, xlsx_coverage
    print(f"[2] ingest done: docx kind={docx_row.kind} pptx kind={pptx_row.kind} xlsx kind={xlsx_row.kind}")

    # 3a. docx read_segment
    storage = get_storage()
    seg = await docx_pl.read_segment(
        file_row=docx_row,
        args={"paragraph_start": 1, "paragraph_end": 2},
        storage=storage,
    )
    assert seg.error is None, seg.error
    assert "Marginalia Test Document" in seg.text
    assert seg.extras["paragraph_start"] == 1
    assert seg.extras["paragraph_end"] == 2
    assert seg.extras["total_paragraphs"] >= 5
    print(f"[3a] docx paragraph 1-2 read: {seg.extras['char_count']} chars")

    seg = await docx_pl.read_segment(
        file_row=docx_row,
        args={"pattern": "raft|leader", "context_lines": 1, "max_matches": 5},
        storage=storage,
    )
    assert seg.error is None, seg.error
    assert seg.extras["match_count"] >= 2
    print(f"[3b] docx pattern matches: {seg.extras['match_count']}")

    # invalid paragraph → clamped, not error
    seg = await docx_pl.read_segment(
        file_row=docx_row,
        args={"paragraph_start": 9999},
        storage=storage,
    )
    # 9999 clamps to last paragraph; not an error
    assert seg.error is None, seg.error
    assert seg.extras["paragraph_start"] == seg.extras["total_paragraphs"]
    print("[3c] out-of-range paragraph_start clamps to end")

    # bad type → error
    seg = await docx_pl.read_segment(
        file_row=docx_row,
        args={"paragraph_start": "not-a-number"},
        storage=storage,
    )
    assert seg.error is not None and "integer" in seg.error
    print("[3d] non-integer paragraph_start → clear error")

    # 4a. pptx read_segment
    seg = await pptx_pl.read_segment(
        file_row=pptx_row,
        args={"slide_start": 1, "slide_end": 1},
        storage=storage,
    )
    assert seg.error is None, seg.error
    assert "Consensus Overview" in seg.text
    assert "Raft leader election" in seg.text
    assert "Presenter note mentions quorum." in seg.text
    assert "Latency Table" not in seg.text
    assert seg.extras["slide_start"] == 1
    assert seg.extras["slide_end"] == 1
    print(f"[4a] pptx slide 1 read: {seg.extras['char_count']} chars")

    seg = await pptx_pl.read_segment(
        file_row=pptx_row,
        args={"pattern": "steady|quorum", "context_lines": 1, "max_matches": 5},
        storage=storage,
    )
    assert seg.error is None, seg.error
    assert seg.extras["match_count"] >= 2
    assert "Round | Latency" in seg.text
    print(f"[4b] pptx pattern matches: {seg.extras['match_count']}")

    seg = await pptx_pl.read_segment(
        file_row=pptx_row,
        args={"page_start": 2},
        storage=storage,
    )
    assert seg.error is None, seg.error
    assert "Latency Table" in seg.text
    assert "Round | Latency" in seg.text
    print("[4c] pptx page_start alias reads slide 2")

    # 5a. spreadsheet read_segment
    seg = await xlsx_pl.read_segment(
        file_row=xlsx_row,
        args={"heading": "consensus"},
        storage=storage,
    )
    assert seg.error is None, seg.error
    assert "consensus" in seg.text
    assert "alice" in seg.text
    assert "latency" not in seg.text  # other sheet not included
    print(f"[5a] spreadsheet heading 'consensus' → {seg.extras['char_count']} chars")

    seg = await xlsx_pl.read_segment(
        file_row=xlsx_row,
        args={"pattern": r"alice|bob", "max_matches": 5},
        storage=storage,
    )
    assert seg.error is None, seg.error
    assert seg.extras["match_count"] >= 2
    print(f"[5b] spreadsheet pattern matches: {seg.extras['match_count']}")

    seg = await xlsx_pl.read_segment(
        file_row=xlsx_row,
        args={"heading": "no-such-sheet"},
        storage=storage,
    )
    assert seg.error is not None
    assert "available_sheets" in seg.extras
    assert "consensus" in seg.extras["available_sheets"]
    print(f"[5c] unknown sheet → error + available_sheets={seg.extras['available_sheets']}")

    # default chunk read includes both sheets
    seg = await xlsx_pl.read_segment(
        file_row=xlsx_row,
        args={"max_chars": 10000},
        storage=storage,
    )
    assert seg.error is None
    assert "consensus" in seg.text and "latency" in seg.text
    print("[5d] default chunk read returns both sheets")

    print("\nALL OFFICE PIPELINES E2E CHECKS PASSED")


if __name__ == "__main__":
    asyncio.run(go())
    sys.exit(0)
